import argparse, json
from pathlib import Path

def read_pdf(path):
    from pdfminer.high_level import extract_text
    return extract_text(str(path)) or ''

def read_pptx(path):
    from pptx import Presentation
    txt = []
    for s in Presentation(str(path)).slides:
        for shape in s.shapes:
            if hasattr(shape, 'text') and shape.text:
                txt.append(shape.text)
    return '\n'.join(txt)

def score(text):
    t = text.lower()
    keys = {
        'exec': ['executive summary', 'objective', 'scope'],
        'cost': ['cost tracking', 'meter', 'usage', 'billing', 'chargeback'],
        'tenant': ['multi tenancy', 'tenant', 'isolation', 'namespace', 'vpc'],
        'security': ['hipaa', 'iso', 'gdpr', 'encryption', 'audit', 'kms'],
        'ops': ['monitor', 'alert', 'incident', 'runbook', 'watch', 'watchdog'],
        'costmodel': ['formula', 'pricing', 'rate card', 'reconciliation'],
        'trade': ['tradeoff', 'risk', 'limitation']
    }
    def anyk(k): return any(x in t for x in keys[k])
    pts = {}
    pts['File validity'] = 10
    pts['Executive summary'] = 10 if anyk('exec') else 3
    pts['Cost tracking'] = 15 if anyk('cost') else 5
    pts['Multi-tenancy'] = 15 if anyk('tenant') else 5
    pts['Security'] = 15 if anyk('security') else 5
    pts['Cost model'] = 15 if anyk('costmodel') else 5
    pts['Operations'] = 10 if anyk('ops') else 3
    pts['Tradeoffs'] = 10 if anyk('trade') else 3
    total = sum(pts.values())
    return pts, min(total, 100)

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--doc', required=True)
    ap.add_argument('--out', required=True)
    a = ap.parse_args()
    p = Path(a.doc)
    if p.suffix.lower() == '.pdf':
        txt = read_pdf(p)
    elif p.suffix.lower() == '.pptx':
        txt = read_pptx(p)
    else:
        txt = ''
    sections, total = score(txt)
    result = {
        'status': 'completed',
        'score': total,
        'sections': [{'name': k, 'score': v} for k, v in sections.items()],
        'feedback': ' | '.join(f"{k}:{v}" for k, v in sections.items())
    }
    Path(a.out).write_text(json.dumps(result, indent=2))

if __name__ == '__main__':
    main()
